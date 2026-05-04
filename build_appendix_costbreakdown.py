"""Appendix slide: cost breakdown tables for the three consolidation savings line items."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

UW_RED = RGBColor(0xC5, 0x05, 0x0C)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Emu(12191695)
SLIDE_H = Emu(6858000)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_textbox(left, top, width, height, text, size=11, bold=False, color=DARK_GRAY,
                font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


# ── UW Red title bar ──────────────────────────────────────
title_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, Inches(0.18)
)
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = UW_RED
title_bar.line.fill.background()

# ── Title ─────────────────────────────────────────────────
add_textbox(
    Inches(0.45), Inches(0.30), Inches(13.0), Inches(0.6),
    "Appendix: Cost Breakdown by Opportunity",
    size=28, bold=True, color=DARK_GRAY,
)
add_textbox(
    Inches(0.45), Inches(0.85), Inches(13.0), Inches(0.35),
    "Source data behind the three consolidation savings figures",
    size=12, bold=False, color=MID_GRAY,
)

# ── Layout: three columns ────────────────────────────────
COL_W = Inches(4.15)
COL_GAP = Inches(0.18)
LEFT_MARGIN = Inches(0.45)
COL_TOP = Inches(1.40)
COL_HEIGHT = Inches(5.10)

col_lefts = [
    LEFT_MARGIN,
    LEFT_MARGIN + COL_W + COL_GAP,
    LEFT_MARGIN + 2 * (COL_W + COL_GAP),
]


def draw_column_header(left, headline, sublabel):
    # Headline pill
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, COL_TOP, COL_W, Inches(0.85)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = UW_RED
    pill.line.fill.background()
    pill.adjustments[0] = 0.18
    tf = pill.text_frame
    tf.margin_left = Emu(72000)
    tf.margin_right = Emu(72000)
    tf.margin_top = Emu(36000)
    tf.margin_bottom = Emu(36000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = headline
    r.font.name = "Calibri"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sublabel
    r2.font.name = "Calibri"
    r2.font.size = Pt(11)
    r2.font.bold = False
    r2.font.color.rgb = WHITE


def draw_table(left, top, width, rows, col_widths, header_fill=LIGHT_GRAY_BG,
               header_size=10, body_size=10):
    """rows = list of list-of-strings; first row is header.
    col_widths = list summing to ~width (in EMU).
    """
    n_rows = len(rows)
    n_cols = len(rows[0])
    # Compute row heights (data rows = 0.30", header = 0.32")
    row_heights = [Inches(0.34)] + [Inches(0.30)] * (n_rows - 1)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                       sum(row_heights, Emu(0)))
    tbl = tbl_shape.table
    for i, h in enumerate(row_heights):
        tbl.rows[i].height = h
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = w

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Emu(54000)
            cell.margin_right = Emu(54000)
            cell.margin_top = Emu(18000)
            cell.margin_bottom = Emu(18000)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            # right-align numeric columns (anything starting with $ or digit or - or +)
            stripped = val.lstrip()
            is_num = stripped.startswith("$") or stripped[:1].isdigit() or stripped.startswith("-")
            if ri == 0:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else (PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT)
            r = p.add_run()
            r.text = val
            r.font.name = "Calibri"
            r.font.size = Pt(header_size if ri == 0 else body_size)
            r.font.bold = (ri == 0) or (ri == n_rows - 1)
            r.font.color.rgb = DARK_GRAY if ri != 0 else DARK_GRAY
    return tbl_shape


# ── Column 1: $700k/yr — Per-Ambulance Costs ─────────────
draw_column_header(col_lefts[0], "$716,818 / yr", "Per-Ambulance Operating Cost")

add_textbox(
    col_lefts[0], COL_TOP + Inches(0.95), COL_W, Inches(0.40),
    "Peterson 24/7 ALS Cost Model (per station)",
    size=11, bold=True, color=DARK_GRAY,
)

c1_rows = [
    ["Line Item", "Annual Cost"],
    ["Salaries", "$371,697"],
    ["Benefits", "$178,466"],
    ["Insurance", "$67,500"],
    ["EMS Supplies", "$28,000"],
    ["WRS Pension", "$27,761"],
    ["Overtime", "$24,894"],
    ["Other (Maint/Equip/Admin/Training)", "$18,500"],
    ["TOTAL OPERATING", "$716,818"],
    ["Less: Revenue Offset", "($466,200)"],
    ["NET COST PER STATION", "$250,618"],
]
c1_widths = [Inches(2.65), Inches(1.50)]
draw_table(col_lefts[0], COL_TOP + Inches(1.35), COL_W, c1_rows, c1_widths,
           body_size=10, header_size=10)

add_textbox(
    col_lefts[0], COL_TOP + Inches(4.85), COL_W, Inches(0.40),
    "Revenue = 700 calls × $666 avg collection.",
    size=8, color=MID_GRAY,
)

# ── Column 2: $59-65k/yr — Overnight FT Staffing ─────────
draw_column_header(col_lefts[1], "$59-65k / yr", "Per-Dept Overnight FT Staffing")

add_textbox(
    col_lefts[1], COL_TOP + Inches(0.95), COL_W, Inches(0.40),
    "Peak-Weighted Shift Savings (per department)",
    size=11, bold=True, color=DARK_GRAY,
)

c2_rows = [
    ["Department", "Night Diff", "On-Call", "OT", "Total"],
    ["Waterloo", "$11,163", "$58,400", "$400", "$69,963"],
    ["Ixonia", "$8,372", "$58,400", "$200", "$66,972"],
    ["Jefferson", "$16,745", "$43,800", "$600", "$61,145"],
    ["Johnson Creek", "$8,372", "$43,800", "$300", "$52,472"],
    ["AVERAGE", "$11,163", "$51,100", "$375", "$62,638"],
    ["TOTAL (4 depts)", "$44,652", "$204,400", "$1,500", "$250,553"],
]
c2_widths = [Inches(1.35), Inches(0.75), Inches(0.75), Inches(0.55), Inches(0.75)]
draw_table(col_lefts[1], COL_TOP + Inches(1.35), COL_W, c2_rows, c2_widths,
           body_size=9, header_size=9)

add_textbox(
    col_lefts[1], COL_TOP + Inches(3.60), COL_W, Inches(1.40),
    "Savings come from:\n"
    "  •  Eliminated 10% night-shift differential on FT salary\n"
    "  •  Reduced overnight on-call hours ($10/hr AEMT, $7.50/hr EMR)\n"
    "  •  Lower afternoon overtime callbacks (better peak coverage)\n\n"
    "Model: shift FT crews from 24/7 to 12-hr peak window (09:00–21:00), "
    "which captures 66% of countywide call volume.",
    size=9, color=DARK_GRAY,
)

# (col 2 source noted in footer)

# ── Column 3: $155k/yr — ALS Hub Net Cost ────────────────
draw_column_header(col_lefts[2], "$155,553 / yr", "Net Cost — 3-Phase Implementation")

add_textbox(
    col_lefts[2], COL_TOP + Inches(0.95), COL_W, Inches(0.40),
    "Combined Phase 1 + 2 + 3 Cash Flow",
    size=11, bold=True, color=DARK_GRAY,
)

c3_rows = [
    ["Phase", "Action", "Annual $"],
    ["1", "Regional ALS Hub dispatch protocol", "$0"],
    ["2", "Peak-weighted FT shifts (4 depts)", "($250,553)"],
    ["3", "County roving paramedic (Lake Mills)", "+$95,000"],
    ["NET ANNUAL IMPACT", "", "($155,553)"],
]
c3_widths = [Inches(0.65), Inches(2.40), Inches(1.10)]
draw_table(col_lefts[2], COL_TOP + Inches(1.35), COL_W, c3_rows, c3_widths,
           body_size=10, header_size=10)

# Roving paramedic detail
add_textbox(
    col_lefts[2], COL_TOP + Inches(3.20), COL_W, Inches(0.35),
    "Phase 3 detail — Roving Paramedic",
    size=11, bold=True, color=DARK_GRAY,
)

c3b_rows = [
    ["Component", "Amount"],
    ["Salary + benefits (1 FTE)", "$95,000"],
    ["Population served", "21,834"],
    ["Overnight calls covered", "231 / yr"],
    ["Cost per capita", "$4.35"],
    ["Est. revenue offset (60% × $666)", "($92,307)"],
    ["NET COST AFTER REVENUE", "$2,693"],
]
c3b_widths = [Inches(2.85), Inches(1.30)]
draw_table(col_lefts[2], COL_TOP + Inches(3.55), COL_W, c3b_rows, c3b_widths,
           body_size=10, header_size=10)

# (sources for col 3 placed in shared footer area below)

# ── Footer ───────────────────────────────────────────────
footer_bg = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35)
)
footer_bg.fill.solid()
footer_bg.fill.fore_color.rgb = LIGHT_GRAY_BG
footer_bg.line.fill.background()

add_textbox(
    Inches(0.45), SLIDE_H - Inches(0.32), Inches(13.0), Inches(0.30),
    "Jefferson County EMS — Act  |  Appendix    Sources: 25-1210 JC EMS Workgroup Cost Projection (Peterson); "
    "Staffing_Reallocation_Recommendations.md; Nighttime_Operations_Deep_Dive.md; Waterloo Chief interview (3/11/26).",
    size=8, color=MID_GRAY,
)

out_path = r"C:\Users\patri\OneDrive - UW-Madison\ISYE 450\appendix_cost_breakdown.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
