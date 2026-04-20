"""
Research Documents Database Builder
Extracts and indexes content from PDFs, spreadsheets, presentations, and Word documents
into a normalized SQLite database with full-text search.
"""

import sqlite3
import os
from pathlib import Path
from datetime import datetime
import hashlib

# Document extraction libraries
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document


def create_database(db_path: str) -> sqlite3.Connection:
    """Create normalized SQLite database with FTS5 for full-text search."""
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Main documents table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path TEXT UNIQUE NOT NULL,
            file_name TEXT NOT NULL,
            file_type TEXT NOT NULL,
            file_size INTEGER,
            file_hash TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            indexed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')

    # Content chunks table (normalized - allows multiple chunks per document)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS content_chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            chunk_index INTEGER NOT NULL,
            chunk_type TEXT,  -- 'page', 'sheet', 'slide', 'section'
            chunk_label TEXT, -- 'Page 1', 'Sheet: Sales', 'Slide 3', etc.
            content TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        )
    ''')

    # Metadata table (key-value pairs for flexible metadata)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS metadata (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            key TEXT NOT NULL,
            value TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        )
    ''')

    # FTS5 virtual table for full-text search
    cursor.execute('''
        CREATE VIRTUAL TABLE IF NOT EXISTS content_fts USING fts5(
            content,
            file_name,
            chunk_label,
            content='content_chunks',
            content_rowid='id'
        )
    ''')

    # Triggers to keep FTS in sync
    cursor.execute('''
        CREATE TRIGGER IF NOT EXISTS content_chunks_ai AFTER INSERT ON content_chunks BEGIN
            INSERT INTO content_fts(rowid, content, file_name, chunk_label)
            SELECT NEW.id, NEW.content, d.file_name, NEW.chunk_label
            FROM documents d WHERE d.id = NEW.document_id;
        END
    ''')

    cursor.execute('''
        CREATE TRIGGER IF NOT EXISTS content_chunks_ad AFTER DELETE ON content_chunks BEGIN
            INSERT INTO content_fts(content_fts, rowid, content, file_name, chunk_label)
            VALUES('delete', OLD.id, OLD.content,
                   (SELECT file_name FROM documents WHERE id = OLD.document_id),
                   OLD.chunk_label);
        END
    ''')

    # Indexes for common queries
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_documents_type ON documents(file_type)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_chunks_document ON content_chunks(document_id)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_metadata_document ON metadata(document_id)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_metadata_key ON metadata(key)')

    conn.commit()
    return conn


def get_file_hash(file_path: str) -> str:
    """Calculate MD5 hash of file for deduplication."""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            hasher.update(chunk)
    return hasher.hexdigest()


def extract_pdf(file_path: str) -> list[dict]:
    """Extract text from PDF, returning list of page chunks."""
    chunks = []
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'page',
                    'chunk_label': f'Page {i + 1}',
                    'content': text.strip()
                })
    except Exception as e:
        print(f"  Error extracting PDF {file_path}: {e}")
    return chunks


def extract_xlsx(file_path: str) -> list[dict]:
    """Extract content from Excel spreadsheet, returning list of sheet chunks."""
    chunks = []
    try:
        wb = load_workbook(file_path, data_only=True)
        for i, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                if any(cell is not None for cell in row):
                    row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                    rows.append(row_text)

            if rows:
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'sheet',
                    'chunk_label': f'Sheet: {sheet_name}',
                    'content': '\n'.join(rows)
                })
    except Exception as e:
        print(f"  Error extracting Excel {file_path}: {e}")
    return chunks


def extract_pptx(file_path: str) -> list[dict]:
    """Extract text from PowerPoint presentation, returning list of slide chunks."""
    chunks = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())

            if texts:
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'slide',
                    'chunk_label': f'Slide {i + 1}',
                    'content': '\n'.join(texts)
                })
    except Exception as e:
        print(f"  Error extracting PowerPoint {file_path}: {e}")
    return chunks


def extract_docx(file_path: str) -> list[dict]:
    """Extract text from Word document."""
    chunks = []
    try:
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = '\t'.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        if paragraphs:
            chunks.append({
                'chunk_index': 0,
                'chunk_type': 'document',
                'chunk_label': 'Full Document',
                'content': '\n\n'.join(paragraphs)
            })
    except Exception as e:
        print(f"  Error extracting Word doc {file_path}: {e}")
    return chunks


def extract_txt(file_path: str) -> list[dict]:
    """Extract text from plain text file."""
    chunks = []
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read().strip()

        if content:
            chunks.append({
                'chunk_index': 0,
                'chunk_type': 'text_file',
                'chunk_label': 'Full Document',
                'content': content
            })
    except Exception as e:
        print(f"  Error extracting text file {file_path}: {e}")
    return chunks


def index_document(conn: sqlite3.Connection, file_path: str, base_path: str) -> bool:
    """Index a single document into the database."""
    cursor = conn.cursor()
    path = Path(file_path)
    relative_path = str(path.relative_to(base_path))
    file_type = path.suffix.lower()

    # Check if already indexed
    cursor.execute('SELECT id FROM documents WHERE file_path = ?', (relative_path,))
    if cursor.fetchone():
        print(f"  Skipping (already indexed): {path.name}")
        return False

    # Extract content based on file type
    extractors = {
        '.pdf': extract_pdf,
        '.xlsx': extract_xlsx,
        '.xls': extract_xlsx,
        '.pptx': extract_pptx,
        '.docx': extract_docx,
        '.txt': extract_txt,
    }

    extractor = extractors.get(file_type)
    if not extractor:
        print(f"  Skipping unsupported file type: {path.name}")
        return False

    print(f"  Extracting: {path.name}")
    chunks = extractor(file_path)

    if not chunks:
        print(f"  No content extracted from: {path.name}")
        return False

    # Insert document record
    file_stat = path.stat()
    file_hash = get_file_hash(file_path)

    cursor.execute('''
        INSERT INTO documents (file_path, file_name, file_type, file_size, file_hash)
        VALUES (?, ?, ?, ?, ?)
    ''', (relative_path, path.name, file_type, file_stat.st_size, file_hash))

    doc_id = cursor.lastrowid

    # Insert content chunks
    for chunk in chunks:
        cursor.execute('''
            INSERT INTO content_chunks (document_id, chunk_index, chunk_type, chunk_label, content)
            VALUES (?, ?, ?, ?, ?)
        ''', (doc_id, chunk['chunk_index'], chunk['chunk_type'], chunk['chunk_label'], chunk['content']))

    # Add metadata
    cursor.execute('INSERT INTO metadata (document_id, key, value) VALUES (?, ?, ?)',
                   (doc_id, 'original_path', str(file_path)))
    cursor.execute('INSERT INTO metadata (document_id, key, value) VALUES (?, ?, ?)',
                   (doc_id, 'chunk_count', str(len(chunks))))

    conn.commit()
    return True


def build_database(base_path: str, db_name: str = 'research_documents.db'):
    """Scan directory and build the research documents database."""
    base_path = Path(base_path)
    db_path = base_path / db_name

    print(f"Building database: {db_path}")
    print(f"Scanning: {base_path}\n")

    conn = create_database(str(db_path))

    # Supported extensions
    extensions = {'.pdf', '.xlsx', '.xls', '.pptx', '.docx', '.txt'}

    # Find all supported files
    files_to_process = []
    for ext in extensions:
        files_to_process.extend(base_path.rglob(f'*{ext}'))

    print(f"Found {len(files_to_process)} documents to process\n")

    indexed = 0
    for file_path in sorted(files_to_process):
        if index_document(conn, str(file_path), str(base_path)):
            indexed += 1

    # Print summary
    cursor = conn.cursor()
    cursor.execute('SELECT COUNT(*) FROM documents')
    doc_count = cursor.fetchone()[0]
    cursor.execute('SELECT COUNT(*) FROM content_chunks')
    chunk_count = cursor.fetchone()[0]

    print(f"\n{'='*50}")
    print(f"Database built successfully!")
    print(f"  Total documents: {doc_count}")
    print(f"  Total content chunks: {chunk_count}")
    print(f"  Database location: {db_path}")
    print(f"\nExample queries:")
    print(f"  -- Search all content:")
    print(f"  SELECT * FROM content_fts WHERE content_fts MATCH 'EMS budget';")
    print(f"  ")
    print(f"  -- Get document with chunks:")
    print(f"  SELECT d.file_name, c.chunk_label, c.content")
    print(f"  FROM documents d JOIN content_chunks c ON d.id = c.document_id")
    print(f"  WHERE d.file_type = '.pdf';")

    conn.close()
    return str(db_path)


if __name__ == '__main__':
    import sys

    if len(sys.argv) > 1:
        base_path = sys.argv[1]
    else:
        # Default to current directory
        base_path = os.path.dirname(os.path.abspath(__file__))

    build_database(base_path)
