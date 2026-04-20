"""
Research Documents Database Builder v2
Optimized schema with entity extraction for Jefferson County EMS Study.

New features:
- Document categories based on folder structure and content
- Municipality extraction and linking
- Service provider extraction and linking
- Contract-specific metadata (dates, payment terms)
- Improved full-text search with category filtering
"""

import sqlite3
import os
import re
from pathlib import Path
from datetime import datetime
import hashlib

# Document extraction libraries
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document


# Known municipalities in Jefferson County
MUNICIPALITIES = {
    'aztalan': 'Town of Aztalan',
    'cambridge': 'Village of Cambridge',
    'cold spring': 'Town of Cold Spring',
    'concord': 'Town of Concord',
    'edgerton': 'City of Edgerton',
    'farmington': 'Town of Farmington',
    'fort atkinson': 'City of Fort Atkinson',
    'hebron': 'Town of Hebron',
    'helenville': 'Town of Helenville',
    'ixonia': 'Town of Ixonia',
    'jefferson': 'City of Jefferson',
    'johnson creek': 'Village of Johnson Creek',
    'koshkonong': 'Town of Koshkonong',
    'lake mills': 'City of Lake Mills',
    'milford': 'Town of Milford',
    'oakland': 'Town of Oakland',
    'palmyra': 'Village of Palmyra',
    'rome': 'Town of Rome',
    'sullivan': 'Town of Sullivan',
    'sumner': 'Town of Sumner',
    'waterloo': 'City of Waterloo',
    'watertown': 'City of Watertown',
    'western lakes': 'Western Lakes Fire District',
    'whitewater': 'City of Whitewater',
    'town of jefferson': 'Town of Jefferson',
    'town of lake mills': 'Town of Lake Mills',
}

# Service providers
SERVICE_PROVIDERS = {
    'jefferson fire': 'Jefferson Fire Department',
    'city of jefferson': 'City of Jefferson EMS',
    'fort atkinson': 'Fort Atkinson Fire/EMS',
    'lake mills': 'Lake Mills Fire/EMS',
    'waterloo': 'Waterloo Fire and Rescue',
    'watertown': 'Watertown Fire Department',
    'edgerton fire protection district': 'Edgerton Fire Protection District',
    'efpd': 'Edgerton Fire Protection District',
    'johnson creek': 'Johnson Creek Fire/EMS',
    'ryan bros': 'Ryan Bros Ambulance',
    'western lakes': 'Western Lakes Fire District',
}

# Document categories based on content/path
CATEGORY_PATTERNS = {
    'contract': (r'contract|agreement|iga\b|amendment', 'EMS Service Contract'),
    'call_data': (r'workgroup.*xlsx|call data', 'EMS Call Data'),
    'policy_report': (r'policy forum|pulse report|greater than the sum', 'Policy Report'),
    'background': (r'background|overview|challenge', 'Background Research'),
    'past_project': (r'univercity|final.*deliverable|badger consulting', 'Past UniverCity Project'),
    'county_example': (r'portage|bayfield|countywide.*ems', 'County EMS Example'),
    'presentation': (r'\.pptx$|presentation', 'Presentation'),
    'meeting_notes': (r'meeting', 'Meeting Notes'),
    'reference': (r'levy.*limit|faq|contacts|resources', 'Reference Document'),
}


def create_database(db_path: str) -> sqlite3.Connection:
    """Create optimized SQLite database with entity tables."""
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # Enable foreign keys
    cursor.execute('PRAGMA foreign_keys = ON')

    # Categories table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS categories (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT UNIQUE NOT NULL,
            description TEXT
        )
    ''')

    # Municipalities table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS municipalities (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT UNIQUE NOT NULL,
            short_name TEXT,
            municipality_type TEXT  -- 'city', 'town', 'village', 'district'
        )
    ''')

    # Service providers table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS service_providers (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT UNIQUE NOT NULL,
            short_name TEXT,
            provider_type TEXT  -- 'fire', 'ems', 'combined', 'private'
        )
    ''')

    # Main documents table (enhanced)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_path TEXT UNIQUE NOT NULL,
            file_name TEXT NOT NULL,
            file_type TEXT NOT NULL,
            file_size INTEGER,
            file_hash TEXT,
            category_id INTEGER,
            folder_path TEXT,
            contract_start_date TEXT,
            contract_end_date TEXT,
            payment_model TEXT, -- e.g., 'Per Capita', 'Equalized Value', 'Fixed'
            payment_rate TEXT, -- e.g., '$31.00/capita (2024)', '85/15 split'
            capital_funding_model TEXT, -- e.g., '85/15 Split', 'Bundled', 'Not Specified'
            has_per_call_fees BOOLEAN,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            indexed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (category_id) REFERENCES categories(id)
        )
    ''')

    # Content chunks table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS content_chunks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            chunk_index INTEGER NOT NULL,
            chunk_type TEXT,
            chunk_label TEXT,
            content TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        )
    ''')

    # Document-Municipality relationship (many-to-many)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_municipalities (
            document_id INTEGER NOT NULL,
            municipality_id INTEGER NOT NULL,
            role TEXT,  -- 'service_area', 'contracting_party', 'mentioned'
            PRIMARY KEY (document_id, municipality_id),
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
            FOREIGN KEY (municipality_id) REFERENCES municipalities(id)
        )
    ''')

    # Document-Service Provider relationship (many-to-many)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS document_providers (
            document_id INTEGER NOT NULL,
            provider_id INTEGER NOT NULL,
            role TEXT,  -- 'provider', 'mentioned'
            PRIMARY KEY (document_id, provider_id),
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
            FOREIGN KEY (provider_id) REFERENCES service_providers(id)
        )
    ''')

    # Metadata table (for additional key-value pairs)
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS metadata (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL,
            key TEXT NOT NULL,
            value TEXT,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        )
    ''')

    # FTS5 virtual table for full-text search
    cursor.execute('''
        CREATE VIRTUAL TABLE IF NOT EXISTS content_fts USING fts5(
            content,
            file_name,
            chunk_label,
            content='content_chunks',
            content_rowid='id'
        )
    ''')

    # Triggers to keep FTS in sync
    cursor.execute('''
        CREATE TRIGGER IF NOT EXISTS content_chunks_ai AFTER INSERT ON content_chunks BEGIN
            INSERT INTO content_fts(rowid, content, file_name, chunk_label)
            SELECT NEW.id, NEW.content, d.file_name, NEW.chunk_label
            FROM documents d WHERE d.id = NEW.document_id;
        END
    ''')

    cursor.execute('''
        CREATE TRIGGER IF NOT EXISTS content_chunks_ad AFTER DELETE ON content_chunks BEGIN
            INSERT INTO content_fts(content_fts, rowid, content, file_name, chunk_label)
            VALUES('delete', OLD.id, OLD.content,
                   (SELECT file_name FROM documents WHERE id = OLD.document_id),
                   OLD.chunk_label);
        END
    ''')

    # Indexes
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_documents_type ON documents(file_type)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_documents_category ON documents(category_id)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_chunks_document ON content_chunks(document_id)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_doc_muni ON document_municipalities(municipality_id)')
    cursor.execute('CREATE INDEX IF NOT EXISTS idx_doc_provider ON document_providers(provider_id)')

    conn.commit()
    return conn


def get_or_create_category(cursor, name: str, description: str = None) -> int:
    """Get existing category or create new one."""
    cursor.execute('SELECT id FROM categories WHERE name = ?', (name,))
    row = cursor.fetchone()
    if row:
        return row[0]
    cursor.execute('INSERT INTO categories (name, description) VALUES (?, ?)', (name, description))
    return cursor.lastrowid


def get_or_create_municipality(cursor, name: str) -> int:
    """Get existing municipality or create new one."""
    cursor.execute('SELECT id FROM municipalities WHERE name = ?', (name,))
    row = cursor.fetchone()
    if row:
        return row[0]

    # Determine type
    muni_type = 'unknown'
    name_lower = name.lower()
    if 'city of' in name_lower:
        muni_type = 'city'
    elif 'town of' in name_lower:
        muni_type = 'town'
    elif 'village of' in name_lower:
        muni_type = 'village'
    elif 'district' in name_lower:
        muni_type = 'district'

    short_name = name.replace('Town of ', '').replace('City of ', '').replace('Village of ', '')
    cursor.execute('INSERT INTO municipalities (name, short_name, municipality_type) VALUES (?, ?, ?)',
                   (name, short_name, muni_type))
    return cursor.lastrowid


def get_or_create_provider(cursor, name: str) -> int:
    """Get existing service provider or create new one."""
    cursor.execute('SELECT id FROM service_providers WHERE name = ?', (name,))
    row = cursor.fetchone()
    if row:
        return row[0]

    provider_type = 'combined'
    if 'ambulance' in name.lower():
        provider_type = 'ems'
    elif 'fire' in name.lower() and 'ems' not in name.lower():
        provider_type = 'fire'

    short_name = name.replace(' Fire Department', '').replace(' Fire/EMS', '').replace(' EMS', '')
    cursor.execute('INSERT INTO service_providers (name, short_name, provider_type) VALUES (?, ?, ?)',
                   (name, short_name, provider_type))
    return cursor.lastrowid


def detect_category(file_path: str, file_name: str, content: str = '') -> str:
    """Detect document category based on path, name, and content."""
    combined = f"{file_path} {file_name} {content[:1000]}".lower()

    for key, (pattern, category_name) in CATEGORY_PATTERNS.items():
        if re.search(pattern, combined, re.IGNORECASE):
            return category_name

    return 'General Document'


def extract_municipalities(content: str) -> list[str]:
    """Extract municipality names from content."""
    found = []
    content_lower = content.lower()

    for key, full_name in MUNICIPALITIES.items():
        # Look for the key or full name
        if key in content_lower or full_name.lower() in content_lower:
            if full_name not in found:
                found.append(full_name)

    return found


def extract_providers(content: str) -> list[str]:
    """Extract service provider names from content."""
    found = []
    content_lower = content.lower()

    for key, full_name in SERVICE_PROVIDERS.items():
        if key in content_lower:
            if full_name not in found:
                found.append(full_name)

    return found


def extract_contract_dates(content: str) -> tuple[str, str]:
    """Extract contract start and end dates from content."""
    start_date = None
    end_date = None

    # Look for date patterns like "January 1, 2023 to December 31, 2027"
    date_pattern = r'(\w+ \d+,? \d{4})\s*(?:to|through|-)\s*(\w+ \d+,? \d{4})'
    match = re.search(date_pattern, content, re.IGNORECASE)
    if match:
        start_date = match.group(1)
        end_date = match.group(2)

    # Also look for year ranges like "2023-2027"
    if not start_date:
        year_pattern = r'(\d{4})\s*[-–]\s*(\d{4})'
        match = re.search(year_pattern, content)
        if match:
            start_date = match.group(1)
            end_date = match.group(2)

    return start_date, end_date


def extract_payment_terms(content: str) -> tuple[str, str]:
    """Extract high-level payment model and specific rate from content."""
    model = None
    rate = None
    content_lower = content.lower()

    if 'per capita' in content_lower:
        model = 'Per Capita'
        # Try to find a specific rate
        rate_match = re.search(r'(\$\d+\.\d{2}\s*/\s*per capita)', content, re.IGNORECASE)
        if rate_match:
            rate = rate_match.group(1)

    elif 'equalized value' in content_lower or 'equalized improvement value' in content_lower:
        model = 'Equalized Value'
        # Try to find a split
        rate_match = re.search(r'(\d+)\s*%\s*of the vehicle', content, re.IGNORECASE)
        if rate_match:
            rate = f"~{rate_match.group(1)}% Town contribution"

    # You can add more specific regex here for other models

    return model, rate


def extract_financial_details(content: str) -> tuple[str, bool]:
    """Extracts capital funding models and presence of per-call fees."""
    capital_model = 'Not Specified'
    per_call_fees = False
    content_lower = content.lower()

    if '85% of the vehicle' in content_lower or '15% of the actual cost' in content_lower:
        capital_model = '85/15 Split'
    elif 'net operating cost' in content_lower:
        capital_model = 'Bundled in Operating Cost'

    if 'hook-up fee' in content_lower or 'out-of-pocket" cost of consumable' in content_lower:
        per_call_fees = True

    return capital_model, per_call_fees

def get_file_hash(file_path: str) -> str:
    """Calculate MD5 hash of file."""
    hasher = hashlib.md5()
    with open(file_path, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            hasher.update(chunk)
    return hasher.hexdigest()


def extract_pdf(file_path: str) -> list[dict]:
    """Extract text from PDF."""
    chunks = []
    try:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'page',
                    'chunk_label': f'Page {i + 1}',
                    'content': text.strip()
                })
    except Exception as e:
        print(f"  Error extracting PDF {file_path}: {e}")
    return chunks


def extract_xlsx(file_path: str) -> list[dict]:
    """Extract content from Excel spreadsheet."""
    chunks = []
    try:
        wb = load_workbook(file_path, data_only=True)
        for i, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                    rows.append(row_text)

            if rows:
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'sheet',
                    'chunk_label': f'Sheet: {sheet_name}',
                    'content': '\n'.join(rows)
                })
    except Exception as e:
        print(f"  Error extracting Excel {file_path}: {e}")
    return chunks


def extract_pptx(file_path: str) -> list[dict]:
    """Extract text from PowerPoint."""
    chunks = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())

            if texts:
                chunks.append({
                    'chunk_index': i,
                    'chunk_type': 'slide',
                    'chunk_label': f'Slide {i + 1}',
                    'content': '\n'.join(texts)
                })
    except Exception as e:
        print(f"  Error extracting PowerPoint {file_path}: {e}")
    return chunks


def extract_docx(file_path: str) -> list[dict]:
    """Extract text from Word document."""
    chunks = []
    try:
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        for table in doc.tables:
            for row in table.rows:
                row_text = '\t'.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        if paragraphs:
            chunks.append({
                'chunk_index': 0,
                'chunk_type': 'document',
                'chunk_label': 'Full Document',
                'content': '\n\n'.join(paragraphs)
            })
    except Exception as e:
        print(f"  Error extracting Word doc {file_path}: {e}")
    return chunks


def extract_txt(file_path: str) -> list[dict]:
    """Extract text from plain text file."""
    chunks = []
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read().strip()

        if content:
            chunks.append({
                'chunk_index': 0,
                'chunk_type': 'text_file',
                'chunk_label': 'Full Document',
                'content': content
            })
    except Exception as e:
        print(f"  Error extracting text file {file_path}: {e}")
    return chunks


def index_document(conn: sqlite3.Connection, file_path: str, base_path: str) -> bool:
    """Index a document with entity extraction."""
    cursor = conn.cursor()
    path = Path(file_path)
    relative_path = str(path.relative_to(base_path))
    file_type = path.suffix.lower()
    folder_path = str(path.parent.relative_to(base_path)) if path.parent != base_path else ''

    # Check if already indexed
    cursor.execute('SELECT id FROM documents WHERE file_path = ?', (relative_path,))
    if cursor.fetchone():
        print(f"  Skipping (already indexed): {path.name}")
        return False

    # Extract content
    extractors = {
        '.pdf': extract_pdf,
        '.xlsx': extract_xlsx,
        '.xls': extract_xlsx,
        '.pptx': extract_pptx,
        '.docx': extract_docx,
        '.txt': extract_txt,
    }

    extractor = extractors.get(file_type)
    if not extractor:
        print(f"  Skipping unsupported file type: {path.name}")
        return False

    print(f"  Extracting: {path.name}")
    chunks = extractor(file_path)

    if not chunks:
        print(f"  No content extracted from: {path.name}")
        return False

    # Combine all content for analysis
    full_content = '\n'.join(chunk['content'] for chunk in chunks)

    # Detect category
    category_name = detect_category(relative_path, path.name, full_content)
    category_id = get_or_create_category(cursor, category_name)

    # Extract contract dates if applicable
    start_date, end_date = None, None
    payment_model, payment_rate = None, None
    if 'contract' in category_name.lower():
        start_date, end_date = extract_contract_dates(full_content)
        payment_model, payment_rate = extract_payment_terms(full_content)

    # Insert document
    file_stat = path.stat()
    file_hash = get_file_hash(file_path)

    cursor.execute('''
        INSERT INTO documents (file_path, file_name, file_type, file_size, file_hash, category_id,
                               folder_path, contract_start_date, contract_end_date, payment_model, payment_rate)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
    ''', (relative_path, path.name, file_type, file_stat.st_size, file_hash,
          category_id, folder_path, start_date, end_date, payment_model, payment_rate))

    doc_id = cursor.lastrowid

    # Insert content chunks
    for chunk in chunks:
        cursor.execute('''
            INSERT INTO content_chunks (document_id, chunk_index, chunk_type, chunk_label, content)
            VALUES (?, ?, ?, ?, ?)
        ''', (doc_id, chunk['chunk_index'], chunk['chunk_type'], chunk['chunk_label'], chunk['content']))

    # Extract and link municipalities
    municipalities = extract_municipalities(full_content)
    for muni_name in municipalities:
        muni_id = get_or_create_municipality(cursor, muni_name)
        role = 'contracting_party' if 'contract' in category_name.lower() else 'mentioned'
        try:
            cursor.execute('INSERT INTO document_municipalities (document_id, municipality_id, role) VALUES (?, ?, ?)',
                           (doc_id, muni_id, role))
        except sqlite3.IntegrityError:
            pass  # Already exists

    # Extract and link service providers
    providers = extract_providers(full_content)
    for provider_name in providers:
        provider_id = get_or_create_provider(cursor, provider_name)
        role = 'provider' if 'contract' in category_name.lower() else 'mentioned'
        try:
            cursor.execute('INSERT INTO document_providers (document_id, provider_id, role) VALUES (?, ?, ?)',
                           (doc_id, provider_id, role))
        except sqlite3.IntegrityError:
            pass

    conn.commit()
    return True


def build_database(base_path: str, db_name: str = 'research_documents.db'):
    """Build the optimized research documents database."""
    base_path = Path(base_path)
    db_path = base_path / db_name

    # Remove old database
    if db_path.exists():
        os.remove(db_path)

    print(f"Building optimized database: {db_path}")
    print(f"Scanning: {base_path}\n")

    conn = create_database(str(db_path))

    # Supported extensions
    extensions = {'.pdf', '.xlsx', '.xls', '.pptx', '.docx', '.txt'}

    # Find all supported files
    files_to_process = []
    for ext in extensions:
        files_to_process.extend(base_path.rglob(f'*{ext}'))

    print(f"Found {len(files_to_process)} documents to process\n")

    indexed = 0
    for file_path in sorted(files_to_process):
        if index_document(conn, str(file_path), str(base_path)):
            indexed += 1

    # Print summary
    cursor = conn.cursor()

    print(f"\n{'='*60}")
    print("DATABASE BUILT SUCCESSFULLY!")
    print('='*60)

    cursor.execute('SELECT COUNT(*) FROM documents')
    print(f"\nTotal documents: {cursor.fetchone()[0]}")

    cursor.execute('SELECT COUNT(*) FROM content_chunks')
    print(f"Total content chunks: {cursor.fetchone()[0]}")

    print("\n--- Documents by Category ---")
    cursor.execute('''
        SELECT c.name, COUNT(d.id)
        FROM categories c
        LEFT JOIN documents d ON c.id = d.category_id
        GROUP BY c.id ORDER BY COUNT(d.id) DESC
    ''')
    for row in cursor.fetchall():
        print(f"  {row[0]}: {row[1]}")

    print("\n--- Municipalities Found ---")
    cursor.execute('SELECT name, municipality_type FROM municipalities ORDER BY name')
    for row in cursor.fetchall():
        print(f"  {row[0]} ({row[1]})")

    print("\n--- Service Providers Found ---")
    cursor.execute('SELECT name, provider_type FROM service_providers ORDER BY name')
    for row in cursor.fetchall():
        print(f"  {row[0]} ({row[1]})")

    print(f"\nDatabase location: {db_path}")

    print("\n--- Example Queries ---")
    print("""
-- Find all contracts for a specific municipality:
SELECT d.file_name, c.name as category, d.contract_start_date, d.contract_end_date
FROM documents d
JOIN categories c ON d.category_id = c.id
JOIN document_municipalities dm ON d.id = dm.document_id
JOIN municipalities m ON dm.municipality_id = m.id
WHERE m.short_name = 'Jefferson';

-- Search content within a category:
SELECT d.file_name, cc.chunk_label, substr(cc.content, 1, 200)
FROM content_fts fts
JOIN content_chunks cc ON fts.rowid = cc.id
JOIN documents d ON cc.document_id = d.id
JOIN categories cat ON d.category_id = cat.id
WHERE content_fts MATCH 'levy' AND cat.name = 'EMS Service Contract';

-- Get all call data by municipality:
SELECT m.short_name, d.file_name
FROM documents d
JOIN document_municipalities dm ON d.id = dm.document_id
JOIN municipalities m ON dm.municipality_id = m.id
JOIN categories c ON d.category_id = c.id
WHERE c.name = 'EMS Call Data'
ORDER BY m.short_name;
""")

    conn.close()
    return str(db_path)


if __name__ == '__main__':
    import sys

    if len(sys.argv) > 1:
        base_path = sys.argv[1]
    else:
        base_path = os.path.dirname(os.path.abspath(__file__))

    build_database(base_path)
