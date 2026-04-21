"""Assemble the Shark Tank pitch deck as a native-PowerPoint .pptx.

Only slides 1 (map) and 2 (bar chart) and 4-left (bar chart) use embedded
PNGs — those are genuine visualizations. Slides 3, 4-right, 5, 6 are built
from native PowerPoint shapes so text stays editable.

Run: python build_pitch_deck.py
Output: pitch_deck_jefferson_ems.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).parent

# ============ Palette ============
NAVY        = RGBColor(0x11, 0x11, 0x11)
RED         = RGBColor(0xB2, 0x22, 0x22)
RED_BG      = RGBColor(0xFF, 0xF5, 0xF5)
BLUE        = RGBColor(0x2C, 0x7F, 0xB8)
BLUE_BG     = RGBColor(0xF0, 0xF8, 0xFF)
GREEN       = RGBColor(0x22, 0x8B, 0x22)
GREEN_BG    = RGBColor(0xF5, 0xFF, 0xF5)
GRAY        = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT  = RGBColor(0x77, 0x77, 0x77)
ORANGE      = RGBColor(0xD9, 0x5F, 0x0E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER     = RGBColor(0xCC, 0xCC, 0xCC)

# ============ Canvas 16:9 ============
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]


# ============ Helpers ============
def textbox(slide, left, top, width, height, text, *,
            size=18, bold=False, italic=False, color=NAVY,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            font="Calibri", line_spacing=1.15):
    """Single-paragraph text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def multiline_textbox(slide, left, top, width, height, lines, *,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                      font="Calibri"):
    """Multi-paragraph text box. `lines` = list of dicts with keys
    text, size, bold, italic, color, space_after (Pt)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.05)

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line.get("space_after"):
            p.space_after = Pt(line["space_after"])
        run = p.add_run()
        run.text = line["text"]
        run.font.size = Pt(line.get("size", 14))
        run.font.bold = line.get("bold", False)
        run.font.italic = line.get("italic", False)
        run.font.color.rgb = line.get("color", NAVY)
        run.font.name = font
    return tb


def rounded_rect(slide, left, top, width, height, *,
                 fill=WHITE, line_color=GRAY_LIGHT, line_width=1.0,
                 corner=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    # No shadow
    sh_elem = shp.shadow
    # Remove text frame default text; we overlay with separate textboxes
    shp.text_frame.text = ""
    return shp


def add_footer(slide, text):
    textbox(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
            text, size=9, italic=True, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, *, color=GRAY, width=3.0):
    """Native horizontal/vertical arrow (connector)."""
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    # Add arrow head via XML tweak
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    # Remove existing tail decoration
    for tag in ("a:headEnd", "a:tailEnd"):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    from lxml import etree
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return conn


# =========================================================
# SLIDE 1 — The Hook
# =========================================================
s1 = prs.slides.add_slide(BLANK)

# Map
s1.shapes.add_picture(
    str(ROOT / "pitch_slide1_anchor_map.png"),
    Inches(0.25), Inches(0.9),
    width=Inches(12.8), height=Inches(6.1),
)

# Top banner: title + team strip
textbox(s1, Inches(0.5), Inches(0.15), Inches(9), Inches(0.5),
        "ISyE 450 Capstone — Jefferson County EMS",
        size=20, bold=True, color=NAVY)
textbox(s1, Inches(9.5), Inches(0.15), Inches(3.5), Inches(0.5),
        "Team: [Your names here]",
        size=14, italic=True, color=GRAY, align=PP_ALIGN.RIGHT)

add_footer(s1,
    "Data: CY2024 NFIRS (n=8,396 EMS calls, Jefferson geography) · "
    "P-Median optimization · OpenRouteService road-network isochrones")


# =========================================================
# SLIDE 2 — The Bottleneck
# =========================================================
s2 = prs.slides.add_slide(BLANK)

textbox(s2, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
        "The Bottleneck: 12 Half-Empty Systems",
        size=34, bold=True, color=NAVY)
textbox(s2, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5),
        "Each of Jefferson County's 12 departments funds its own backup ambulance — "
        "used 10–15% of the time, yet still leaving 41 coverage gaps per year.",
        size=15, italic=True, color=GRAY)

# Chart left
s2.shapes.add_picture(
    str(ROOT / "secondary_demand_by_dept_jeffco.png"),
    Inches(0.4), Inches(1.7),
    width=Inches(8.3), height=Inches(5.2),
)

# Native "key numbers" panel right
rounded_rect(s2, Inches(9.0), Inches(1.9), Inches(3.9), Inches(5.0),
             fill=RED_BG, line_color=RED, line_width=2.0, corner=0.06)

textbox(s2, Inches(9.2), Inches(2.05), Inches(3.5), Inches(0.5),
        "THE KEY NUMBERS",
        size=15, bold=True, color=RED)

multiline_textbox(s2, Inches(9.2), Inches(2.55), Inches(3.5), Inches(4.2), [
    {"text": "21",      "size": 36, "bold": True, "color": RED,  "space_after": 2},
    {"text": "all-busy events / year\nin Ixonia (≈ once / 2.5 weeks)",
     "size": 11, "italic": True, "color": GRAY, "space_after": 12},
    {"text": "12",      "size": 36, "bold": True, "color": RED,  "space_after": 2},
    {"text": "all-busy events / year\nin Waterloo",
     "size": 11, "italic": True, "color": GRAY, "space_after": 12},
    {"text": "41",      "size": 36, "bold": True, "color": RED,  "space_after": 2},
    {"text": "countywide — NO ambulance\nwas available for that 911 call",
     "size": 11, "italic": True, "color": GRAY, "space_after": 0},
])

add_footer(s2, "Source: CY2024 NFIRS concurrent-call analysis · "
               "Jefferson County geography only")


# =========================================================
# SLIDE 3 — The IE Intervention (native shapes, no image)
# =========================================================
s3 = prs.slides.add_slide(BLANK)

textbox(s3, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
        "The IE Intervention: Hybrid Local + Regional",
        size=34, bold=True, color=NAVY)
textbox(s3, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5),
        "Primaries stay local. Secondaries go regional. "
        "Dispatch routes whichever unit reaches the patient fastest.",
        size=15, italic=True, color=BLUE)

# ---- LEFT column: TODAY ----
LX, LY, LW = Inches(0.5), Inches(1.8), Inches(5.8)
rounded_rect(s3, LX, LY, LW, Inches(5.0),
             fill=RED_BG, line_color=RED, line_width=2.0)

textbox(s3, LX, Inches(1.9), LW, Inches(0.5),
        "TODAY",
        size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
textbox(s3, LX, Inches(2.4), LW, Inches(0.4),
        "12 independent departments",
        size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# 5 small dept "cards" in a row
dept_y = Inches(2.95)
dept_h = Inches(2.3)
card_w_in = 1.05
gap_in = 0.05
total_w = 5 * card_w_in + 4 * gap_in  # 5.45
start_x_in = 0.5 + (5.8 - total_w) / 2
for i in range(5):
    cx = Inches(start_x_in + i * (card_w_in + gap_in))
    rounded_rect(s3, cx, dept_y, Inches(card_w_in), dept_h,
                 fill=WHITE, line_color=RED, line_width=1.2)
    textbox(s3, cx, Inches(3.05), Inches(card_w_in), Inches(0.3),
            f"Dept {chr(65+i)}", size=10, bold=True, color=RED,
            align=PP_ALIGN.CENTER)
    # Primary: solid box
    rounded_rect(s3, Inches(start_x_in + i * (card_w_in + gap_in) + 0.25),
                 Inches(3.4), Inches(0.55), Inches(0.55),
                 fill=RED, line_color=RED, line_width=0.5)
    textbox(s3, cx, Inches(4.0), Inches(card_w_in), Inches(0.3),
            "Primary", size=9, color=NAVY, align=PP_ALIGN.CENTER)
    # Secondary: outlined box (idle)
    rounded_rect(s3, Inches(start_x_in + i * (card_w_in + gap_in) + 0.25),
                 Inches(4.3), Inches(0.55), Inches(0.55),
                 fill=WHITE, line_color=GRAY_LIGHT, line_width=1.0)
    textbox(s3, cx, Inches(4.9), Inches(card_w_in), Inches(0.3),
            "Backup", size=9, italic=True, color=GRAY_LIGHT,
            align=PP_ALIGN.CENTER)

# Pain-point strip at bottom of LEFT
textbox(s3, LX, Inches(5.9), LW, Inches(0.8),
        "Each department funds its own backup.\n"
        "Still 41 all-busy events per year.",
        size=13, bold=True, italic=True, color=RED, align=PP_ALIGN.CENTER,
        line_spacing=1.25)

# ---- Arrow divider ----
add_arrow(s3, Inches(6.45), Inches(4.2), Inches(7.05), Inches(4.2),
          color=NAVY, width=4.0)

# ---- RIGHT column: PROPOSED ----
RX, RY, RW = Inches(7.15), Inches(1.8), Inches(5.7)
rounded_rect(s3, RX, RY, RW, Inches(5.0),
             fill=BLUE_BG, line_color=BLUE, line_width=2.0)

textbox(s3, RX, Inches(1.9), RW, Inches(0.5),
        "PROPOSED",
        size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
textbox(s3, RX, Inches(2.4), RW, Inches(0.4),
        "Local primaries + 4 shared regional secondaries",
        size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Primaries row (smaller, all blue)
prim_y_in = 3.0
prim_card_w = 0.95
prim_total = 5 * prim_card_w + 4 * gap_in
prim_start_x = 7.15 + (5.7 - prim_total) / 2
for i in range(5):
    cx = Inches(prim_start_x + i * (prim_card_w + gap_in))
    rounded_rect(s3, cx, Inches(prim_y_in), Inches(prim_card_w), Inches(0.9),
                 fill=WHITE, line_color=BLUE, line_width=1.0)
    textbox(s3, cx, Inches(prim_y_in + 0.08), Inches(prim_card_w), Inches(0.3),
            f"Dept {chr(65+i)}", size=9, bold=True, color=BLUE,
            align=PP_ALIGN.CENTER)
    rounded_rect(s3, Inches(prim_start_x + i * (prim_card_w + gap_in) + 0.3),
                 Inches(prim_y_in + 0.40), Inches(0.35), Inches(0.35),
                 fill=BLUE, line_color=BLUE, line_width=0.5)

textbox(s3, RX, Inches(4.05), RW, Inches(0.35),
        "Primaries stay local — autonomy preserved",
        size=11, italic=True, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Regional stations row — 4 red stars
star_y_in = 4.5
star_count = 4
star_spacing = 1.1
star_total = (star_count - 1) * star_spacing
star_start_x = 7.15 + (5.7 - star_total) / 2 - 0.15
for i in range(star_count):
    sx = Inches(star_start_x + i * star_spacing)
    shp = s3.shapes.add_shape(MSO_SHAPE.STAR_5_POINT, sx,
                              Inches(star_y_in), Inches(0.7), Inches(0.7))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RED
    shp.line.color.rgb = WHITE
    shp.line.width = Pt(1.5)
    shp.text_frame.text = ""
    textbox(s3, Inches(star_start_x + i * star_spacing - 0.15),
            Inches(star_y_in + 0.75), Inches(1.0), Inches(0.35),
            f"SEC-{i+1}", size=11, bold=True, color=RED,
            align=PP_ALIGN.CENTER)

textbox(s3, RX, Inches(5.95), RW, Inches(0.75),
        "4 shared regional secondaries.\n"
        "Dispatch routes the closest unit — regardless of town lines.",
        size=12, bold=True, italic=True, color=BLUE, align=PP_ALIGN.CENTER,
        line_spacing=1.25)

add_footer(s3, "Methodology: P-Median facility location (Amazon logistics) · "
               "Erlang-C queuing (call-center staffing) · "
               "OpenRouteService drive-time validation")


# =========================================================
# SLIDE 4 — The Data Evidence (hybrid: chart + native cards)
# =========================================================
s4 = prs.slides.add_slide(BLANK)

textbox(s4, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7),
        "From Fragmented Backup to Regional Resilience",
        size=32, bold=True, color=NAVY)

# ---- LEFT: BEFORE bar chart ----
textbox(s4, Inches(0.5), Inches(1.1), Inches(6.5), Inches(0.5),
        "BEFORE — Fragmented coverage",
        size=22, bold=True, color=RED)
textbox(s4, Inches(0.5), Inches(1.55), Inches(6.5), Inches(0.5),
        "41 times in 2024, a caller's primary ambulance was busy AND "
        "the department had no backup.",
        size=12, italic=True, color=GRAY, line_spacing=1.15)

s4.shapes.add_picture(
    str(ROOT / "pitch_slide4_bars_only.png"),
    Inches(0.3), Inches(2.2),
    width=Inches(6.9), height=Inches(4.6),
)

# ---- RIGHT: AFTER native KPI cards ----
textbox(s4, Inches(7.4), Inches(1.1), Inches(5.6), Inches(0.5),
        "AFTER — K=4 Regional Network",
        size=22, bold=True, color=BLUE)
textbox(s4, Inches(7.4), Inches(1.55), Inches(5.6), Inches(0.5),
        "4 regional secondary stations overlay existing primaries.",
        size=12, italic=True, color=GRAY)

# Three native KPI cards, stacked
kpis = [
    {"label": "Demand covered ≤14 min",
     "today": "Today: 32%", "after": "K=4: 75%", "delta": "+43 pts",
     "color": BLUE, "bg": BLUE_BG},
    {"label": "Avg secondary response time",
     "today": "Today: 16.6 min", "after": "K=4: 10.8 min", "delta": "−5.8 min",
     "color": BLUE, "bg": BLUE_BG},
    {"label": "All-busy events closed by overflow",
     "today": "Today: 41 / yr", "after": "K=4: ≤5 / yr", "delta": "−36 events",
     "color": GREEN, "bg": GREEN_BG},
]

card_top_in = 2.2
card_h_in = 1.45
card_gap_in = 0.1
for i, k in enumerate(kpis):
    y = Inches(card_top_in + i * (card_h_in + card_gap_in))
    rounded_rect(s4, Inches(7.4), y, Inches(5.6), Inches(card_h_in),
                 fill=k["bg"], line_color=k["color"], line_width=2.0,
                 corner=0.1)
    # Label
    textbox(s4, Inches(7.55), y + Inches(0.12), Inches(5.3), Inches(0.4),
            k["label"], size=13, bold=True, color=GRAY)
    # Today (left)
    textbox(s4, Inches(7.55), y + Inches(0.55), Inches(2.0), Inches(0.5),
            k["today"], size=14, color=GRAY_LIGHT)
    # Arrow
    add_arrow(s4, Inches(9.6), y + Inches(0.75),
              Inches(10.0), y + Inches(0.75),
              color=k["color"], width=3.0)
    # After (right)
    textbox(s4, Inches(10.15), y + Inches(0.55), Inches(2.2), Inches(0.5),
            k["after"], size=16, bold=True, color=k["color"])
    # Delta chip (top-right)
    textbox(s4, Inches(7.55), y + Inches(1.0), Inches(5.2), Inches(0.4),
            k["delta"], size=16, bold=True, color=k["color"],
            align=PP_ALIGN.RIGHT)

add_footer(s4, "Before: K=2 MCLP T=10 baseline for current fragmented "
               "mutual-aid · After: K=4 P-Median optimization, total demand")


# =========================================================
# SLIDE 5 — The Bottom Line (native cards)
# =========================================================
s5 = prs.slides.add_slide(BLANK)

textbox(s5, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
        "The Bottom Line — For the Patient",
        size=38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
textbox(s5, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5),
        "In the language that matters to a 911 caller",
        size=15, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Three big hero cards
hero = [
    {"label": "TIME",     "headline": "−5.8 MIN",
     "sub": "Faster response during\nconcurrent calls",
     "detail": "16.6 min → 10.8 min.\nIn cardiac arrest, this is\nthe difference that matters.",
     "color": RED, "bg": RED_BG},
    {"label": "QUALITY",  "headline": "1",
     "sub": "Unified medical protocol,\ncounty-wide",
     "detail": "Today: 12 separate protocols.\nTomorrow: every patient gets\nthe same standard of care.",
     "color": BLUE, "bg": BLUE_BG},
    {"label": "COVERAGE", "headline": "41 → 5",
     "sub": "All-busy events closed\nper year",
     "detail": "No Jefferson County resident\nwaits for mutual aid from\nanother county.",
     "color": GREEN, "bg": GREEN_BG},
]

card_w_in = 3.9
card_h_in = 5.2
gap_in = 0.25
total_w = 3 * card_w_in + 2 * gap_in  # 12.0
start_x_in = (13.333 - total_w) / 2
card_top = Inches(1.9)
for i, c in enumerate(hero):
    cx = Inches(start_x_in + i * (card_w_in + gap_in))
    rounded_rect(s5, cx, card_top, Inches(card_w_in), Inches(card_h_in),
                 fill=c["bg"], line_color=c["color"], line_width=3.0,
                 corner=0.08)
    # Label pill
    textbox(s5, cx, card_top + Inches(0.3), Inches(card_w_in), Inches(0.5),
            c["label"], size=18, bold=True, color=c["color"],
            align=PP_ALIGN.CENTER)
    # Headline
    textbox(s5, cx, card_top + Inches(1.0), Inches(card_w_in), Inches(1.3),
            c["headline"], size=56, bold=True, color=c["color"],
            align=PP_ALIGN.CENTER)
    # Sub
    textbox(s5, cx, card_top + Inches(2.5), Inches(card_w_in), Inches(1.0),
            c["sub"], size=15, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, line_spacing=1.2)
    # Detail
    textbox(s5, cx, card_top + Inches(3.7), Inches(card_w_in), Inches(1.3),
            c["detail"], size=12, italic=True, color=GRAY,
            align=PP_ALIGN.CENTER, line_spacing=1.25)

textbox(s5, Inches(0.5), Inches(7.2), Inches(12.3), Inches(0.3),
        "We're not asking for funding. We're asking for partnership as the "
        "Jefferson County EMS Working Group moves into implementation.",
        size=12, bold=True, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 5b — The Business Case (Peterson-anchored cost model)
# =========================================================
s5b = prs.slides.add_slide(BLANK)

textbox(s5b, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
        "The Business Case: Same Spend, Dramatically Better Care",
        size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
textbox(s5b, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.5),
        "Built on Chief Peterson's 24/7 ALS single-unit cost model "
        "(6 FTEs, 24/48 schedule, 700 calls/yr)",
        size=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# ---- TODAY panel (left) — the diagnostic baseline ----
rounded_rect(s5b, Inches(0.5), Inches(1.9), Inches(4.0), Inches(4.5),
             fill=RED_BG, line_color=RED, line_width=2.0)

textbox(s5b, Inches(0.5), Inches(2.05), Inches(4.0), Inches(0.5),
        "TODAY — Fragmented",
        size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

textbox(s5b, Inches(0.65), Inches(2.6), Inches(3.7), Inches(0.4),
        "Current distributed overhead",
        size=12, color=GRAY)

textbox(s5b, Inches(0.65), Inches(2.95), Inches(3.7), Inches(0.9),
        "$2.36M",
        size=48, bold=True, color=RED)

textbox(s5b, Inches(0.65), Inches(3.85), Inches(3.7), Inches(0.4),
        "per year, across 7 departments",
        size=12, italic=True, color=GRAY)

# Divider line
line = s5b.shapes.add_connector(2, Inches(0.7), Inches(4.35),
                                 Inches(4.3), Inches(4.35))
line.line.color.rgb = DIVIDER
line.line.width = Pt(1.0)

multiline_textbox(s5b, Inches(0.7), Inches(4.5), Inches(3.6), Inches(1.8), [
    {"text": "What that buys today:", "size": 11, "bold": True,
     "color": GRAY, "space_after": 6},
    {"text": "• 12 independent backup ambulances",
     "size": 11, "color": NAVY, "space_after": 3},
    {"text": "• 10–15% utilization per unit",
     "size": 11, "color": NAVY, "space_after": 3},
    {"text": "• 41 all-busy gaps per year",
     "size": 11, "color": RED, "bold": True, "space_after": 3},
    {"text": "• Only 32% of concurrent demand covered ≤14 min",
     "size": 11, "color": NAVY, "space_after": 0},
])

# ---- ARROW ----
add_arrow(s5b, Inches(4.7), Inches(4.0), Inches(5.1), Inches(4.0),
          color=NAVY, width=4.0)
textbox(s5b, Inches(4.4), Inches(3.5), Inches(0.9), Inches(0.4),
        "Same\nbudget",
        size=11, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER,
        line_spacing=1.1)

# ---- TOMORROW panel (right) — 3 scenarios, Peterson-anchored ----
rounded_rect(s5b, Inches(5.3), Inches(1.9), Inches(7.55), Inches(4.5),
             fill=BLUE_BG, line_color=BLUE, line_width=2.0)

textbox(s5b, Inches(5.3), Inches(2.05), Inches(7.55), Inches(0.5),
        "TOMORROW — Regional (K=4)",
        size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# 3 scenario cards across, scaled from Peterson's $250K-net single-unit model
scenarios = [
    {"name": "Scenario A", "mode": "All 4 stations 24/7 ALS",
     "op": "$2.87M", "rev": "$1.86M", "net": "$1.00M",
     "fte": "29 FTE",
     "note": "Maximum coverage.\nHigher cost; delivers ≤14-min\ncoverage to 75% of demand."},
    {"name": "Scenario B", "mode": "All 4 peak-only (8 AM–8 PM)",
     "op": "$2.06M", "rev": "$1.21M", "net": "$851K",
     "fte": "19 FTE",
     "note": "Peak-hour focus — covers\n90% of call volume during\nthe 12-hour daytime window."},
    {"name": "Scenario C", "mode": "Hybrid: 1×24/7 + 3×peak",
     "op": "$2.33M", "rev": "$1.43M", "net": "$902K",
     "fte": "22 FTE",
     "note": "Recommended balance.\n24/7 anchor + peak-only\ncoverage where demand spikes."},
]

card_w_in = 2.35
card_gap_in = 0.12
card_top = Inches(2.65)
card_h_in = 3.55
start_x_in = 5.3 + 0.12
for i, sc in enumerate(scenarios):
    cx = Inches(start_x_in + i * (card_w_in + card_gap_in))
    is_rec = sc["name"] == "Scenario C"
    border = GREEN if is_rec else BLUE
    bw = 2.5 if is_rec else 1.2
    rounded_rect(s5b, cx, card_top, Inches(card_w_in), Inches(card_h_in),
                 fill=WHITE, line_color=border, line_width=bw)
    # Scenario name
    textbox(s5b, cx, card_top + Inches(0.1), Inches(card_w_in), Inches(0.35),
            sc["name"], size=13, bold=True, color=border, align=PP_ALIGN.CENTER)
    # Mode
    textbox(s5b, cx, card_top + Inches(0.45), Inches(card_w_in), Inches(0.5),
            sc["mode"], size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER,
            line_spacing=1.15)
    # Net cost (big)
    textbox(s5b, cx, card_top + Inches(1.05), Inches(card_w_in), Inches(0.55),
            sc["net"], size=26, bold=True, color=border, align=PP_ALIGN.CENTER)
    textbox(s5b, cx, card_top + Inches(1.60), Inches(card_w_in), Inches(0.3),
            "net / year", size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    # Op / Rev breakdown
    multiline_textbox(s5b, cx + Inches(0.05), card_top + Inches(1.95),
                      Inches(card_w_in - 0.1), Inches(0.7), [
        {"text": f"Operating: {sc['op']}", "size": 10, "color": GRAY,
         "space_after": 2},
        {"text": f"Revenue:  {sc['rev']}", "size": 10, "color": GRAY,
         "space_after": 2},
        {"text": sc["fte"], "size": 10, "bold": True, "color": NAVY,
         "space_after": 0},
    ], align=PP_ALIGN.CENTER)
    # Note
    textbox(s5b, cx + Inches(0.1), card_top + Inches(2.7),
            Inches(card_w_in - 0.2), Inches(0.8), sc["note"],
            size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER,
            line_spacing=1.2)
    # "Recommended" chip for Scenario C
    if is_rec:
        textbox(s5b, cx, card_top - Inches(0.28), Inches(card_w_in), Inches(0.3),
                "★ RECOMMENDED", size=10, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)

# Bottom takeaway strip
textbox(s5b, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.6),
        "The hybrid regional network costs ≈ $900K net — within the envelope "
        "the county already spends on fragmented backup.",
        size=14, bold=True, italic=True, color=NAVY, align=PP_ALIGN.CENTER,
        line_spacing=1.2)

add_footer(s5b,
    "Source: Chief Peterson 24/7 ALS cost projection (6 FTEs, 700 calls/yr) "
    "scaled to K=4 · Current distributed overhead estimated from 7-dept staffing data")


# =========================================================
# SLIDE 6 — Thank You + Call to Action
# =========================================================
s6 = prs.slides.add_slide(BLANK)

textbox(s6, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5),
        "Thank You.",
        size=72, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

textbox(s6, Inches(1.5), Inches(3.1), Inches(10.3), Inches(1.0),
        "Partnership and feedback as the Jefferson County EMS "
        "Working Group moves into implementation.",
        size=20, italic=True, color=BLUE, align=PP_ALIGN.CENTER,
        line_spacing=1.3)

textbox(s6, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.5),
        "ISyE 450 Capstone Team",
        size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

textbox(s6, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.6),
        "[Team member 1]    [Team member 2]    [Team member 3]    [Team member 4]",
        size=14, color=GRAY, align=PP_ALIGN.CENTER)

textbox(s6, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.5),
        "Questions? We have backup slides on sensitivity, staffing scenarios, "
        "and dispatch protocols.",
        size=13, italic=True, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

add_footer(s6, "University of Wisconsin–Madison · ISyE 450 · April 2026")


# ============ Save ============
out = ROOT / "pitch_deck_jefferson_ems.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
