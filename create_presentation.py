"""
Generate Final Presentation PowerPoint for Jefferson County EMS Study
ISyE 450 Senior Design — April 2026
Clean version: large readable text, tables with big fonts, NO heatmaps
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = r"C:\Users\patri\OneDrive - UW-Madison\ISYE 450"

# ── Colors ──
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT    = RGBColor(0xC0, 0x39, 0x2B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xF2, 0xF2, 0xF2)
MED_GRAY  = RGBColor(0x7F, 0x8C, 0x8D)
BLACK     = RGBColor(0x33, 0x33, 0x33)
GOLD      = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──
def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    return tf

def add_bullets(slide, left, top, width, height, bullets, font_size=18,
                color=BLACK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return tf

def add_image_safe(slide, path, left, top, width=None, height=None):
    full = os.path.join(BASE, path)
    if os.path.exists(full):
        kwargs = {}
        if width: kwargs['width'] = Inches(width)
        if height: kwargs['height'] = Inches(height)
        slide.shapes.add_picture(full, Inches(left), Inches(top), **kwargs)
        return True
    return False

def section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    add_text_box(slide, 0.8, 2.0, 11.7, 1.2, title, font_size=44, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, 0.8, 3.7, 11.7, 1.0, subtitle, font_size=22, color=LIGHT_GRAY)
    return slide

def content_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    add_text_box(slide, 0.6, 0.15, 12, 0.75, title, font_size=30, bold=True, color=WHITE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(2.5), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    return slide

def add_table(slide, left, top, width, height, data, col_widths=None,
              header_color=NAVY, font_size=16):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = BLACK
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
    return table

# ============================================================
# 1. TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_text_box(slide, 1.0, 1.8, 11.3, 1.5,
             "Jefferson County EMS", font_size=56, bold=True, color=WHITE)
add_text_box(slide, 1.0, 3.2, 11.3, 1.0,
             "Improving EMS Care Through Regional Coordination",
             font_size=30, color=LIGHT_GRAY)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.3), Inches(3.5), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_text_box(slide, 1.0, 4.6, 11.3, 0.5,
             "ISyE 450 Senior Design  |  April 2026", font_size=20, color=MED_GRAY)
add_text_box(slide, 1.0, 5.2, 11.3, 0.5,
             "Prepared for the Jefferson County EMS Working Group", font_size=18, color=MED_GRAY)

# ============================================================
# 2. AGENDA
# ============================================================
slide = content_slide("Agenda")
items = [
    "The Challenge: Why Jefferson County EMS Needs Attention",
    "Current System: 14,853 Calls Across 12 Providers",
    "Secondary Demand: When Primary Ambulances Are Busy",
    "Goal 1: Regional Secondary Ambulance Network",
    "Optimal Station Placement & Cost Scenarios",
    "Goal 2: Peak Staffing Investigation",
    "Where County-Funded EMTs Help Most",
    "Recommendations & Next Steps",
]
for i, item in enumerate(items):
    items[i] = f"{i+1}.  {item}"
add_bullets(slide, 1.0, 1.5, 11, 5.5, items, font_size=22)

# ============================================================
# 3. THE CHALLENGE (section)
# ============================================================
section_slide("The Challenge",
    "How can Jefferson County improve EMS response quality for its residents?")

# ============================================================
# 4. CURRENT SYSTEM
# ============================================================
slide = content_slide("Current System: Fragmented Municipal Model")
bullets = [
    "12 independent EMS providers across 576 sq mi  (pop. 86,855)",
    "14,853 total EMS calls in 2024 - growing with aging population",
    "No county EMS coordinator, no centralized dispatch, no shared billing",
    "Each department maintains its own ambulances, staffing, and contracts",
    "Municipalities at or near levy caps - EMS competes with other services",
    "WI Stat. 66.0602: County EMS levy is EXEMPT from levy limits (unused)",
]
add_bullets(slide, 0.6, 1.3, 7.5, 5.5, bullets, font_size=18)

call_data = [
    ["Provider", "2024 Calls"],
    ["Western Lakes", "5,633"],
    ["Edgerton/Lakeside", "2,138"],
    ["Watertown", "2,012"],
    ["Fort Atkinson", "1,616"],
    ["Jefferson", "1,457"],
    ["Waterloo", "520"],
    ["Lake Mills", "518"],
    ["Johnson Creek", "487"],
    ["Ixonia", "289"],
    ["Other (3 depts)", "183"],
    ["TOTAL", "14,853"],
]
add_table(slide, 8.5, 1.3, 4.3, 5.5, call_data, col_widths=[2.6, 1.7], font_size=15)

# ============================================================
# 5. COUNTY COMPARISON
# ============================================================
slide = content_slide("How Other WI Counties Do It")
comp = [
    ["", "Jefferson Co.", "Portage Co.", "Bayfield Co."],
    ["Population", "86,855", "70,521", "~15,000"],
    ["Model", "Municipal\n(no county role)", "Hybrid\n(county contracts)", "County-supported\n(stipends)"],
    ["County EMS Levy", "None", "Yes ($4M+)", "Yes ($458K)"],
    ["Ambulances", "18 across 7 depts", "7 (4 staffed 24/7)", "Volunteer"],
    ["EMS Coordinator", "None", "Yes (since 2009)", "Proposed 2025"],
    ["Medical Direction", "Per department", "1 MD countywide", "Not specified"],
]
add_table(slide, 0.6, 1.3, 12, 4.8, comp, col_widths=[2.5, 3.2, 3.2, 3.1], font_size=16)

add_text_box(slide, 0.6, 6.3, 12, 0.6,
    "Jefferson County runs 2.5x more ambulances than Portage for a similar population, "
    "but with no county-level coordination.",
    font_size=17, bold=True, color=ACCENT)

# ============================================================
# 6. SECONDARY DEMAND (section)
# ============================================================
section_slide("Secondary Demand",
    "When a department's primary ambulance is already on a call, the next patient waits.")

# ============================================================
# 7. SECONDARY DEMAND TABLE
# ============================================================
slide = content_slide("How Often Are All Ambulances Busy?")

sec = [
    ["Department", "Ambulances", "EMS Calls", "Secondary\nEvents", "All-Busy\nEvents"],
    ["Edgerton", "2", "2,035", "768 (37.7%)", "161 (7.9%)"],
    ["Watertown", "3", "1,947", "656 (33.7%)", "16 (0.8%)"],
    ["Whitewater", "2", "1,448", "421 (29.1%)", "63 (4.4%)"],
    ["Fort Atkinson", "2", "1,621", "209 (12.9%)", "1 (0.1%)"],
    ["Waterloo", "2", "403", "93 (23.1%)", "12 (3.0%)"],
    ["Johnson Creek", "1", "454", "59 (13.0%)", "7 (1.5%)"],
    ["Ixonia", "1", "260", "32 (12.3%)", "32 (12.3%)"],
]
add_table(slide, 0.6, 1.3, 12, 4.5, sec, col_widths=[2.2, 1.5, 1.5, 2.8, 2.8], font_size=17)

add_text_box(slide, 0.6, 6.0, 12, 0.5,
    "County-wide: 2,244 secondary demand events per year.  "
    "Peak hours 09:00-19:00 (2.9x the overnight rate).",
    font_size=18, bold=True, color=DARK_BLUE)
add_text_box(slide, 0.6, 6.7, 12, 0.5,
    "Source: CY2024 NFIRS data. \"Secondary event\" = another call active while this dept's ambulance was on a call.",
    font_size=13, color=MED_GRAY)

# ============================================================
# 8. SECONDARY DEMAND BAR CHART
# ============================================================
slide = content_slide("Secondary Demand by Department")
if not add_image_safe(slide, "secondary_demand_by_dept.png", 1.5, 1.3, width=10.5):
    add_text_box(slide, 2, 3, 9, 2, "[Chart: secondary_demand_by_dept.png]", font_size=24, color=MED_GRAY)

# ============================================================
# 9. ERLANG-C TABLE
# ============================================================
slide = content_slide("Probability a Patient Has to Wait (Erlang-C)")

erlang = [
    ["Department", "Ambulances", "P(wait)\nAll-Day", "P(wait)\nPeak Hrs", "What This Means"],
    ["Edgerton", "2", "2.5%", "4.5%", "1 in 22 peak calls waits"],
    ["Ixonia", "1", "3.5%", "4.7%", "Single-unit: any overlap = gap"],
    ["Whitewater", "2", "0.9%", "1.4%", "Moderate congestion"],
    ["Waterloo", "2", "0.2%", "0.5%", "Manageable"],
    ["Watertown", "3", "0.07%", "0.15%", "Well-served by 3 units"],
]
add_table(slide, 0.6, 1.3, 12, 3.8, erlang, col_widths=[2.0, 1.6, 1.6, 1.6, 5.2], font_size=17)

add_bullets(slide, 0.6, 5.5, 12, 1.8, [
    "P(wait) = probability ALL ambulances are busy when a new call comes in",
    "Edgerton and Ixonia are the two departments where patients are most likely to wait",
], font_size=17, color=BLACK)

# ============================================================
# 10. GOAL 1 SECTION
# ============================================================
section_slide("Goal 1: Regional Secondary Ambulance Network",
    "Consolidate backup ambulance capacity into a shared, optimized countywide network")

# ============================================================
# 11. CONCEPT
# ============================================================
slide = content_slide("How It Would Work")
left_bullets = [
    "TODAY: Each department independently maintains backup ambulances",
    "  - Staffed by part-time workers at 10-15% utilization",
    "  - $2.36M/year total across 7 departments",
    "",
    "PROPOSED: Dedicated regional secondary stations",
    "  - When primary ambulance is busy, call goes to nearest regional unit",
    "  - Professional full-time crews (reliable, trained)",
    "  - Shared resource serving ALL municipalities",
    "  - Placed at locations that minimize response time county-wide",
]
add_bullets(slide, 0.6, 1.3, 6.5, 5.5, left_bullets, font_size=17)

add_text_box(slide, 7.8, 1.3, 5, 0.5, "What This Achieves", font_size=22, bold=True, color=NAVY)
right_bullets = [
    "Higher utilization of each ambulance",
    "Consistent full-time coverage (not on-call PT)",
    "Fewer total ambulances to maintain",
    "Better overflow dispatch and communication",
]
add_bullets(slide, 7.8, 2.1, 5, 4, right_bullets, font_size=18, color=DARK_BLUE)

# ============================================================
# 12. HOW MANY STATIONS
# ============================================================
slide = content_slide("How Many Stations Do We Need?")

stations = [
    ["# Stations", "14-min Coverage", "Avg Response Time"],
    ["2", "66.8%", "11.7 min"],
    ["3", "86.1%", "10.6 min"],
    ["4", "91.2%", "9.7 min"],
    ["5", "96.0%", "7.9 min"],
]
add_table(slide, 0.6, 1.3, 5.5, 3.0, stations, col_widths=[1.5, 2.0, 2.0], font_size=18)

add_bullets(slide, 0.6, 4.5, 5.5, 2.5, [
    "Going from 2 to 3 stations: +19 points of coverage",
    "Going from 3 to 4 stations: only +5 points",
    "3 stations is the practical sweet spot",
], font_size=17)

# Diminishing returns chart on right
if not add_image_safe(slide, "secondary_network_diminishing_returns.png", 6.5, 1.2, width=6.3):
    add_text_box(slide, 7, 3, 5, 2, "[Chart: diminishing returns]", font_size=20, color=MED_GRAY)

# ============================================================
# 13. K=3 MAP
# ============================================================
slide = content_slide("Where to Put the 3 Stations")
if not add_image_safe(slide, "secondary_network_map_K3.png", 1.5, 1.2, width=10):
    add_text_box(slide, 3, 3, 7, 2, "[Map: K=3 station placement]", font_size=24, color=MED_GRAY)
add_text_box(slide, 1.0, 6.6, 11, 0.6,
    "Three zones:  North (Watertown area)  |  Central (Jefferson / Fort Atkinson)  |  South (Whitewater / Edgerton)",
    font_size=18, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# ============================================================
# 14. COST SCENARIOS
# ============================================================
slide = content_slide("What Would It Cost?")

cost = [
    ["Scenario", "Hours", "Net Cost/yr", "FTE Needed"],
    ["A: All 3 stations 24/7", "24/7", "$751,854", "21.6"],
    ["B: All 3 peak-only (08-20)", "12 hr/day", "$638,546", "14.4"],
    ["C: Hybrid (1 x 24/7 + 2 x peak)", "Mixed", "$676,315", "16.8"],
]
add_table(slide, 0.6, 1.3, 12, 2.8, cost, col_widths=[4.0, 2.0, 2.5, 2.0], font_size=18)

add_text_box(slide, 0.6, 4.3, 6, 0.5, "Based on Chief Peterson's Cost Model:", font_size=18, bold=True, color=NAVY)
add_bullets(slide, 0.6, 4.9, 6, 2.5, [
    "Operating cost per 24/7 ALS station: $716,818/yr",
    "Revenue per station: $466,200/yr  (700 calls x $666 avg)",
    "Net cost per station: $250,618/yr",
], font_size=17)

# Waterfall chart on right
if not add_image_safe(slide, "staffing_waterfall.png", 7.0, 4.0, width=5.8):
    add_text_box(slide, 7.5, 4.5, 5, 2, "[Chart: cost waterfall]", font_size=20, color=MED_GRAY)

# ============================================================
# 15. CURRENT vs CONSOLIDATED
# ============================================================
slide = content_slide("Current Distributed Cost vs. Consolidated Network")

if not add_image_safe(slide, "current_vs_consolidated.png", 0.5, 1.2, width=7.5):
    add_text_box(slide, 1, 3, 6, 2, "[Chart: current vs consolidated]", font_size=20, color=MED_GRAY)

add_bullets(slide, 8.5, 1.5, 4.5, 5, [
    "Current: $2.36M/year spread across 7 departments for independent backup capacity",
    "",
    "Consolidated 3-station network: $676K-$752K/yr for 86% of secondary demand",
    "",
    "~51 part-time positions become 16.8 full-time (professional, dedicated crews)",
    "",
    "This is a diagnostic finding - the Working Group decides what to do with it",
], font_size=17)

# ============================================================
# 16. GOAL 2 SECTION
# ============================================================
section_slide("Goal 2: Peak Staffing Investigation",
    "If the county provides EMTs, where do they improve patient care the most?")

# ============================================================
# 17. WHEN DO CALLS PEAK
# ============================================================
slide = content_slide("When Do EMS Calls Happen?")

if not add_image_safe(slide, "peak_staffing_optimal_shift.png", 0.5, 1.2, width=7.5):
    add_text_box(slide, 1, 3, 6, 2, "[Chart: hourly call profile]", font_size=20, color=MED_GRAY)

add_bullets(slide, 8.5, 1.5, 4.3, 5, [
    "Peak hours: 09:00 - 19:00",
    "3.5x more calls at peak vs overnight",
    "Best 8-hour window: 08:00-16:00 (43% of all calls)",
    "Best 12-hour window: 07:00-19:00 (65% of all calls)",
    "Most depts staff uniformly 24/7 despite this pattern",
], font_size=18)

# ============================================================
# 18. OPTIMAL EMT PLACEMENT
# ============================================================
slide = content_slide("If the County Funds EMTs, Where Should They Go?")

emt = [
    ["# of EMTs", "Best Assignment", "Why"],
    ["1st", "Edgerton Day (08-20)", "6x more impact than any other option"],
    ["2nd", "Whitewater Day (08-20)", "2nd highest peak-hour congestion"],
    ["3rd", "Edgerton Night (20-08)", "Night coverage for busiest dept"],
    ["4th", "Ixonia Day (08-20)", "Single ambulance - any overlap = zero coverage"],
    ["5th", "Whitewater Night (20-08)", "Completes 24/7 for 2nd busiest"],
]
add_table(slide, 0.6, 1.3, 12, 4.0, emt, col_widths=[1.5, 3.5, 7.0], font_size=18)

add_bullets(slide, 0.6, 5.5, 12, 1.8, [
    "Edgerton: 5.6 calls/day with 61-min avg duration = highest utilization in the county",
    "Ixonia ranks #4 despite low volume because its single ambulance means ANY overlap is a total gap",
    "Method: Erlang-C queueing model, marginal P(wait) reduction per department per shift",
], font_size=16)

# ============================================================
# 19. STAFFING MISMATCH
# ============================================================
slide = content_slide("The Staffing Mismatch")

if not add_image_safe(slide, "peak_staffing_overstaffing.png", 0.5, 1.2, width=6.5):
    add_text_box(slide, 1, 3, 5, 2, "[Chart: staffing mismatch]", font_size=20, color=MED_GRAY)

add_bullets(slide, 7.5, 1.5, 5.3, 5, [
    "Most departments staff the same 24/7",
    "But 65% of calls happen 07:00-19:00",
    "Edgerton 09:00-15:00 consistently exceeds statistical control limits (SPC 2-sigma)",
    "Watertown and Whitewater also show regular peak-hour strain",
    "Matching staffing to demand patterns improves care when it matters most",
], font_size=18)

# ============================================================
# 20. ISyE TOOLS
# ============================================================
slide = content_slide("ISyE 450 Methods We Used")

tools = [
    ["Method", "What We Did With It"],
    ["Erlang-C Queueing", "Calculated probability a patient waits, for each dept and hour"],
    ["Integer Programming\n(MCLP / P-Median)", "Found optimal station locations from 60 candidates"],
    ["Pareto Analysis", "Balanced coverage vs. cost vs. number of stations"],
    ["SPC Control Charts", "Flagged hours where demand exceeds normal limits"],
    ["Sweep-Line Algorithm", "Detected 2,244 concurrent call events in 13,800+ records"],
    ["PDCA Framework", "Structured our work: Plan phase done, Act phase next"],
]
add_table(slide, 0.6, 1.3, 12, 5.0, tools, col_widths=[3.5, 8.5], font_size=17)

# ============================================================
# 21. RECOMMENDATIONS
# ============================================================
slide = content_slide("Recommendations")

add_text_box(slide, 0.6, 1.3, 6, 0.5, "Near-Term", font_size=24, bold=True, color=NAVY)
add_bullets(slide, 0.6, 1.9, 6, 3.5, [
    "1. Hire a county EMS coordinator",
    "2. First county EMT at Edgerton daytime (08-20)",
    "3. Address Ixonia single-ambulance vulnerability",
    "4. Unify medical direction across departments",
], font_size=18)

add_text_box(slide, 0.6, 5.0, 6, 0.5, "Longer-Term", font_size=24, bold=True, color=NAVY)
add_bullets(slide, 0.6, 5.6, 6, 1.5, [
    "5. Explore county EMS levy (exempt from levy cap)",
    "6. Pilot 3-station regional secondary network",
], font_size=18)

# Key numbers on right
add_text_box(slide, 7.5, 1.3, 5.3, 0.5, "Key Numbers", font_size=24, bold=True, color=NAVY)
nums = [
    ["", ""],
    ["Total EMS Calls (2024)", "14,853"],
    ["Secondary Demand Events/yr", "2,244"],
    ["Optimal Secondary Stations", "3 (86% coverage)"],
    ["Edgerton All-Busy Events", "161/year"],
    ["Ixonia All-Busy Rate", "12.3%"],
    ["Current Overhead (7 depts)", "$2.36M/year"],
    ["Consolidated Network Cost", "$676-752K/yr"],
    ["Best County EMT Placement", "Edgerton Day"],
    ["Best Shift Window", "08:00-16:00"],
]
add_table(slide, 7.3, 1.9, 5.5, 5.0, nums, col_widths=[3.0, 2.5], font_size=16)

# ============================================================
# 22. NEXT STEPS
# ============================================================
slide = content_slide("Next Steps")

add_bullets(slide, 0.8, 1.4, 11.5, 5.5, [
    "Present these findings to the full EMS Working Group",
    "",
    "We provide the data - fire chiefs and county officials decide what to act on",
    "",
    "Possible Act-phase items (if the Working Group directs):",
    "   - Implementation plan for regional secondary network",
    "   - Dispatch protocol for secondary ambulance routing",
    "   - PT-to-FT staffing transition roadmap",
    "   - County EMS levy feasibility (MS Business students' scope)",
    "",
    "Interactive dashboard: jefferson-county-ems.fly.dev",
], font_size=20)

# ============================================================
# 23. DATA SOURCES
# ============================================================
slide = content_slide("Data Sources")
src = [
    ["Source", "Period", "Used For"],
    ["14 NFIRS Excel files", "CY2024", "Call volumes, temporal patterns, response times"],
    ["Department budgets", "FY2025", "Expense, revenue, staffing levels"],
    ["Chief Peterson cost model", "Dec 2025", "24/7 ALS station cost projection"],
    ["WI DOA / Census ACS", "2024-25", "Service area populations"],
    ["ORS road-network times", "Cached", "Station-to-block-group drive times"],
    ["Fire chief interviews", "Mar 2026", "Staffing corrections, operations context"],
    ["MABAS Div 118 sheets", "2025", "Ambulance fleet inventory (18 units)"],
    ["Chief Association data", "2024-25", "Billing collections for 9 agencies"],
]
add_table(slide, 0.6, 1.3, 12, 5.0, src, col_widths=[3.0, 1.5, 7.5], font_size=16)

# ============================================================
# 24. THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_text_box(slide, 1.0, 2.2, 11.3, 1.5, "Thank You", font_size=56, bold=True, color=WHITE)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.5), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
add_text_box(slide, 1.0, 4.1, 11.3, 0.6,
             "ISyE 450 Senior Design  |  Jefferson County EMS Study", font_size=22, color=LIGHT_GRAY)
add_text_box(slide, 1.0, 4.8, 11.3, 0.6,
             "Dashboard:  jefferson-county-ems.fly.dev", font_size=20, color=MED_GRAY)
add_text_box(slide, 1.0, 5.5, 11.3, 0.6, "Questions?", font_size=32, color=GOLD)

# ── Save ──
out = os.path.join(BASE, "Jefferson_County_EMS_Final_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
